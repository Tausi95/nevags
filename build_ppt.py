"""NEVAGS DEC Presentation — Professional redesign matching template style"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design tokens ─────────────────────────────────────────
F      = "Calibri"          # font family
DARK   = RGBColor(0x1C,0x2B,0x28)  # near-black heading
GRAY   = RGBColor(0x4B,0x55,0x63)  # body text
LGRAY  = RGBColor(0x9C,0xA3,0xAF)  # captions / metadata
CARD   = RGBColor(0xF3,0xF4,0xF6)  # card background
LCARD  = RGBColor(0xF8,0xFA,0xFC)  # lighter card
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GREEN  = RGBColor(0x1B,0x43,0x32)  # NEVAGS forest green
GREEN2 = RGBColor(0x2D,0x6A,0x4F)  # mid green
GREEN3 = RGBColor(0x40,0x91,0x6C)  # light green
ORANGE = RGBColor(0xE8,0x69,0x0A)  # NEVAGS orange
AMBER  = RGBColor(0xF5,0x9E,0x0B)  # amber
RED    = RGBColor(0xDC,0x26,0x26)

# ── Canvas ────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Primitive helpers ─────────────────────────────────────
def R(sl, x, y, w, h, fill, border=None):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def T(sl, text, x, y, w, h, sz=12, bold=False, color=DARK,
      align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color; r.font.name = F
    return tb

def P(tf, text, sz=11, bold=False, color=GRAY,
      align=PP_ALIGN.LEFT, sp=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    if sp: p.space_before = Pt(sp)
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color; r.font.name = F
    return p

# ── Layout components ──────────────────────────────────────
STRIP_H = Inches(0.50)   # header strip height
CONTENT_Y = Inches(1.80) # where body content starts

def page_header(sl, section, page, invert=False):
    """Thin top strip: company name | section | page number, then two rule lines."""
    bg = GREEN if invert else WHITE
    R(sl, 0, 0, W, STRIP_H, bg)
    lo = WHITE if invert else ORANGE
    lm = WHITE if invert else LGRAY
    T(sl, "NEVAGS ECO BRICK & CONSTRUCTION",
      Inches(0.4), Inches(0.13), Inches(5), Inches(0.26),
      sz=9, bold=True, color=lo)
    T(sl, section.upper(),
      Inches(4.5), Inches(0.13), Inches(5.5), Inches(0.26),
      sz=9, color=lm, align=PP_ALIGN.CENTER)
    T(sl, f"May 2026  ·  {page:02d}",
      Inches(11.2), Inches(0.13), Inches(1.95), Inches(0.26),
      sz=9, color=lm, align=PP_ALIGN.RIGHT)
    R(sl, 0, STRIP_H,            W, Inches(0.022), GREEN)
    R(sl, 0, STRIP_H+Inches(0.022), W, Inches(0.008), ORANGE)

def section_head(sl, label, heading, page, invert=False):
    """Section label + big bold heading + short green rule."""
    page_header(sl, label, page, invert)
    ht = WHITE if invert else DARK
    T(sl, label.upper(),
      Inches(0.5), Inches(0.70), Inches(12), Inches(0.28),
      sz=9, bold=True, color=ORANGE)
    T(sl, heading,
      Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.75),
      sz=34, bold=True, color=ht)
    R(sl, Inches(0.5), Inches(1.68), Inches(1.1), Inches(0.038), GREEN)

def stat_card(sl, x, y, w, h, number, label, sub="",
              num_color=GREEN, bg=CARD):
    """Large-number stat card — clean, minimal."""
    R(sl, x, y, w, h, bg)
    R(sl, x, y, w, Inches(0.038), GREEN)          # top accent
    T(sl, number,
      x+Inches(0.2), y+Inches(0.1), w-Inches(0.4), Inches(0.6),
      sz=30, bold=True, color=num_color)
    T(sl, label,
      x+Inches(0.2), y+Inches(0.66), w-Inches(0.4), Inches(0.26),
      sz=10, bold=True, color=DARK)
    if sub:
        T(sl, sub,
          x+Inches(0.2), y+Inches(0.90), w-Inches(0.4), Inches(0.28),
          sz=8.5, color=LGRAY)

def card(sl, x, y, w, h, title, bullets,
         bg=CARD, accent=GREEN, tc=DARK, bc=GRAY, tsz=11, bsz=10):
    """Content card with left accent bar."""
    R(sl, x, y, w, h, bg)
    R(sl, x, y, Inches(0.038), h, accent)
    T(sl, title,
      x+Inches(0.18), y+Inches(0.14), w-Inches(0.28), Inches(0.32),
      sz=tsz, bold=True, color=tc)
    bx = sl.shapes.add_textbox(
        x+Inches(0.18), y+Inches(0.50), w-Inches(0.28), h-Inches(0.58))
    bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = b; r.font.size = Pt(bsz); r.font.color.rgb = bc
        r.font.name = F

def table(sl, x, y, w, rows, cols, data, cw=None,
          hbg=GREEN, alt=RGBColor(0xF8,0xFA,0xFC)):
    rh = Inches(0.40)
    tbl = sl.shapes.add_table(rows, cols, x, y, w, rh*rows).table
    if cw:
        for i, c in enumerate(cw): tbl.columns[i].width = c
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top  = cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; run = p.add_run()
            val = data[r][c] if r < len(data) and c < len(data[r]) else ""
            run.text = str(val); run.font.name = F
            if r == 0:
                run.font.size = Pt(10); run.font.bold = True
                run.font.color.rgb = WHITE
                f = cell.fill; f.solid(); f.fore_color.rgb = hbg
            else:
                run.font.size = Pt(10); run.font.color.rgb = GRAY
                f = cell.fill; f.solid()
                f.fore_color.rgb = alt if r % 2 == 0 else WHITE

def footer_bar(sl, text, bg=GREEN, tc=WHITE, sz=10):
    R(sl, 0, H-Inches(0.42), W, Inches(0.42), bg)
    T(sl, text, Inches(0.5), H-Inches(0.38), Inches(12.33), Inches(0.32),
      sz=sz, bold=True, color=tc, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)

# Right dark panel
R(sl, Inches(8.3), 0, Inches(5.03), H, GREEN)

# Orange top stripe
R(sl, 0, 0, Inches(8.3), Inches(0.06), ORANGE)

# Left content
T(sl, "DISTRICT EXECUTIVE COMMITTEE",
  Inches(0.55), Inches(0.28), Inches(7.5), Inches(0.28),
  sz=10, bold=True, color=ORANGE)
T(sl, "Mulanje District, Malawi  ·  May 2026",
  Inches(0.55), Inches(0.55), Inches(7.5), Inches(0.28),
  sz=10, color=LGRAY)

# Company name — large
T(sl, "NEVAGS",
  Inches(0.5), Inches(1.2), Inches(7.7), Inches(1.5),
  sz=88, bold=True, color=GREEN)
T(sl, "ECO BRICK & CONSTRUCTION",
  Inches(0.55), Inches(2.65), Inches(7.5), Inches(0.55),
  sz=22, bold=True, color=DARK)
T(sl, '"Building Tomorrow Sustainably"',
  Inches(0.55), Inches(3.22), Inches(7.5), Inches(0.38),
  sz=14, italic=True, color=GRAY)

# Thin rule
R(sl, Inches(0.55), Inches(3.68), Inches(4.5), Inches(0.025), GREEN)

# Stat row
for i, (num, lbl, sub) in enumerate([
    ("K450M+",  "Capital Invested",   "Founder + Atmosfair EUR 175K"),
    ("51",      "Current Staff",      "15 Female · 36 Male"),
    ("K2.07B",  "Revenue Potential",  "Annual MWK at capacity"),
    ("8,000+",  "Jobs in 3 Years",    "Direct + ecosystem"),
]):
    sx = Inches(0.55 + i * 1.92)
    T(sl, num,  sx, Inches(3.82), Inches(1.88), Inches(0.5),
      sz=18, bold=True, color=GREEN)
    T(sl, lbl,  sx, Inches(4.28), Inches(1.88), Inches(0.24),
      sz=8.5, bold=True, color=DARK)
    T(sl, sub,  sx, Inches(4.5),  Inches(1.88), Inches(0.28),
      sz=7.5, color=LGRAY)

# Pioneer note
R(sl, Inches(0.55), Inches(5.1), Inches(7.5), Inches(0.34),
  RGBColor(0xFF,0xF3,0xE0))
T(sl, "One of the first black-owned industrial companies in Mulanje District",
  Inches(0.65), Inches(5.14), Inches(7.3), Inches(0.26),
  sz=9, bold=True, color=ORANGE)

# ── Right panel content ──
T(sl, "Reg. New Vision Anenenji Construction",
  Inches(8.55), Inches(0.25), Inches(4.5), Inches(0.28),
  sz=8.5, color=RGBColor(0x90,0xB8,0xA0))
T(sl, "Reg. No. 46289",
  Inches(8.55), Inches(0.50), Inches(4.5), Inches(0.24),
  sz=8.5, color=RGBColor(0x70,0x98,0x80))

R(sl, Inches(8.55), Inches(1.0), Inches(4.5), Inches(0.025),
  RGBColor(0x40,0x60,0x50))

T(sl, "COMMISSIONING\nCEREMONY",
  Inches(8.55), Inches(1.1), Inches(4.5), Inches(0.7),
  sz=11, bold=True, color=ORANGE)
T(sl, "27 May 2026",
  Inches(8.55), Inches(1.75), Inches(4.5), Inches(0.7),
  sz=38, bold=True, color=WHITE)
T(sl, "Official Factory Opening\nMusewu, Mulanje District",
  Inches(8.55), Inches(2.42), Inches(4.5), Inches(0.5),
  sz=11, color=RGBColor(0xCC,0xDD,0xCC))

R(sl, Inches(8.55), Inches(3.1), Inches(4.5), Inches(0.025),
  RGBColor(0x40,0x60,0x50))

for i, (num, lbl) in enumerate([
    ("VSK Technology",    "Zero-firewood production"),
    ("SDGs 3 5 13 15",    "UN sustainability goals"),
    ("GIZ · Atmosfair",   "International backed"),
    ("Musewu, Mulanje",   "Community-rooted"),
]):
    T(sl, num, Inches(8.55), Inches(3.25 + i*0.72), Inches(4.5), Inches(0.28),
      sz=11, bold=True, color=WHITE)
    T(sl, lbl, Inches(8.55), Inches(3.50 + i*0.72), Inches(4.5), Inches(0.25),
      sz=9, color=RGBColor(0xAA,0xCC,0xBB))

# Bottom bar
R(sl, 0, H-Inches(0.38), W, Inches(0.38), DARK)
T(sl, "Charles Billy Nasala (MD)  ·  +265 888 34 75 75  ·  nasalacharles.b@gmail.com  "
  "·  Chancy Tausi Tsonga (BD)  ·  +265 984 000 366  ·  chancy.tsonga@yahoo.com",
  Inches(0.5), H-Inches(0.34), Inches(12.33), Inches(0.28),
  sz=8.5, color=RGBColor(0xAA,0xBB,0xAA), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Presentation Overview", "What We Will Cover Today", 2)

items = [
    ("01", "Company Background",    "Why NEVAGS exists · the firewood ban · our opportunity"),
    ("02", "Leadership & Products", "Ownership, registration · VSK bricks · biomass briquettes"),
    ("03", "Investment & Market",   "Capital base · returns · cost-per-m² comparison"),
    ("04", "Environment & SDGs",    "Zero firewood · SDG alignment · zero-waste model"),
    ("05", "Employment Plan",       "51 staff → 8,000+ jobs · community moulding programme"),
    ("06", "Budget & Partners",     "3-month operational budget · GIZ, Atmosfair & others"),
    ("07", "Roadmap",               "27 May 2026 commissioning · phased growth to 2029"),
    ("08", "Partnership Ask",       "What NEVAGS needs from Mulanje District"),
]
cols = [items[:4], items[4:]]
for ci, col in enumerate(cols):
    x = Inches(0.5 + ci * 6.45)
    for ri, (num, title, desc) in enumerate(col):
        y = CONTENT_Y + Inches(ri * 1.25)
        R(sl, x, y, Inches(6.2), Inches(1.18), CARD)
        R(sl, x, y, Inches(0.038), Inches(1.18), GREEN)
        T(sl, num,   x+Inches(0.18), y+Inches(0.1),  Inches(0.5), Inches(0.38),
          sz=22, bold=True, color=RGBColor(0xD1,0xD5,0xDB))
        T(sl, title, x+Inches(0.75), y+Inches(0.12), Inches(5.3), Inches(0.35),
          sz=13, bold=True, color=DARK)
        T(sl, desc,  x+Inches(0.75), y+Inches(0.50), Inches(5.3), Inches(0.55),
          sz=9.5, color=GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — COMPANY BACKGROUND
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Company Background", "Why NEVAGS Exists", 3)

# Two-column layout
for ci, (bg_c, accent_c, ttl, items) in enumerate([
    (CARD, GREEN, "THE PROBLEM — MALAWI'S HOUSING & DEFORESTATION CRISIS", [
        "80–90% of Malawi's housing is built with fired clay bricks",
        "For decades, bricks were fired using massive quantities of firewood",
        "Malawi has one of the fastest deforestation rates in Africa",
        "Soil erosion, flooding and biodiversity loss have accelerated",
        "",
        "MALAWI GOVERNMENT BANNED firewood-fired bricks",
        "Critical supply gap — developers, NGOs and govt need a legal alternative",
        "Demand is policy-driven, long-term and nationwide",
    ]),
    (GREEN, ORANGE, "THE SOLUTION — NEVAGS ECO BRICK & CONSTRUCTION", [
        "VSK (Vertical Shaft Kiln) technology fuelled by biomass briquettes",
        "Zero firewood — 100% policy compliant, every day",
        "High-quality bricks at affordable pricing",
        "Located in the heart of Mulanje District",
        "Already employing 51 community members",
        "Backed by GIZ, Atmosfair & academic research partners",
        "",
        '"When we build, let us think that we build forever." — Ruskin',
    ]),
]):
    x = Inches(0.5 + ci * 6.45)
    R(sl, x, CONTENT_Y, Inches(6.2), Inches(5.45), bg_c)
    R(sl, x, CONTENT_Y, Inches(0.038), Inches(5.45), accent_c)
    T(sl, ttl,
      x+Inches(0.18), CONTENT_Y+Inches(0.12), Inches(5.9), Inches(0.32),
      sz=9, bold=True, color=ORANGE if ci==0 else AMBER)
    R(sl, x+Inches(0.18), CONTENT_Y+Inches(0.46), Inches(5.85), Inches(0.018),
      RGBColor(0xCC,0xCC,0xCC) if ci==0 else RGBColor(0x40,0x70,0x50))
    bx = sl.shapes.add_textbox(
        x+Inches(0.18), CONTENT_Y+Inches(0.52), Inches(5.85), Inches(4.7))
    bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
    first = True
    for b in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        if not b:
            p.space_before = Pt(6); r = p.add_run(); r.text = ""; r.font.name = F
            continue
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = b
        r.font.size = Pt(10)
        r.font.name = F
        if ci == 0:
            r.font.color.rgb = RED if "BANNED" in b else GRAY
            r.font.bold = "BANNED" in b or "policy-driven" in b
        else:
            r.font.color.rgb = WHITE if not b.startswith('"') else AMBER
            r.font.italic = b.startswith('"')


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — LEADERSHIP & OWNERSHIP
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Company Ownership", "Leadership & Registration", 4)

# Pioneer banner
R(sl, Inches(0.5), CONTENT_Y, Inches(12.33), Inches(0.38),
  RGBColor(0xFF,0xF7,0xED))
R(sl, Inches(0.5), CONTENT_Y, Inches(0.038), Inches(0.38), ORANGE)
T(sl, "One of the first black-owned industrial companies in Mulanje District "
  "— a pioneer enterprise building community ownership",
  Inches(0.65), CONTENT_Y+Inches(0.06), Inches(12.0), Inches(0.26),
  sz=9.5, bold=True, color=ORANGE)

# Three person/info cards
CARD_Y = CONTENT_Y + Inches(0.5)
CARD_H = Inches(4.85)

for ci, (bg, accent, role, name, lines) in enumerate([
    (GREEN, ORANGE,
     "FOUNDER, OWNER & MANAGING DIRECTOR", "Charles Billy Nasala",
     ["Entrepreneur, industrial innovator & community champion.",
      "Founder of one of Malawi's first VSK brick enterprises.",
      "",
      "T:  +265 888 34 75 75",
      "T:  +265 99 751 0160",
      "E:  nasalacharles.b@gmail.com"]),
    (DARK, ORANGE,
     "FOUNDING ENGINEER & BD MANAGER", "Chancy Tausi Tsonga",
     ["VSK technology specialist & strategic partner.",
      "Leads business development and market expansion.",
      "",
      "T:  +265 984 000 366",
      "WA: +27 764 998 4601",
      "E:  chancy.tsonga@yahoo.com",
      "W:  chancytsonga.com"]),
    (CARD, GREEN,
     "COMPANY REGISTRATION", "New Vision Anenenji Construction",
     ["Registration No. 46289",
      "",
      "Location",
      "Musewu, Mulanje District, Southern Malawi",
      "P.O. Box 90, Mulanje",
      "",
      "Careers / HR",
      "careers.nevags@gmail.com"]),
]):
    x = Inches(0.5 + ci * 4.28)
    R(sl, x, CARD_Y, Inches(4.1), CARD_H, bg)
    R(sl, x, CARD_Y, Inches(0.038), CARD_H, accent)
    T(sl, role,
      x+Inches(0.18), CARD_Y+Inches(0.14), Inches(3.75), Inches(0.38),
      sz=8.5, bold=True, color=ORANGE if ci < 2 else ORANGE)
    name_c = WHITE if ci < 2 else DARK
    T(sl, name,
      x+Inches(0.18), CARD_Y+Inches(0.52), Inches(3.75), Inches(0.55),
      sz=16, bold=True, color=name_c)
    R(sl, x+Inches(0.18), CARD_Y+Inches(1.08), Inches(3.6), Inches(0.015),
      RGBColor(0x40,0x65,0x50) if ci < 2 else RGBColor(0xCC,0xCC,0xCC))
    bx = sl.shapes.add_textbox(
        x+Inches(0.18), CARD_Y+Inches(1.15), Inches(3.7), Inches(3.5))
    bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
    first = True
    for b in lines:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run(); r.text = b; r.font.name = F
        if not b:
            r.font.size = Pt(4); r.font.color.rgb = bg; continue
        if b.startswith(("T:", "E:", "WA:", "W:")):
            r.font.size = Pt(9.5)
            r.font.color.rgb = AMBER if ci < 2 else ORANGE
        elif b in ("Location", "Careers / HR", "Operational Reach"):
            r.font.size = Pt(9.5); r.font.bold = True
            r.font.color.rgb = WHITE if ci < 2 else DARK
        else:
            r.font.size = Pt(9.5)
            r.font.color.rgb = RGBColor(0xCC,0xDD,0xCC) if ci < 2 else GRAY


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — MISSION, VISION & OBJECTIVES
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Mission, Vision & Objectives", "Our Strategic Foundation", 5)

COL_Y = CONTENT_Y
COL_H = Inches(5.42)

for ci, (bg, ttl, sub, body, extras) in enumerate([
    (GREEN, "MISSION", "What We Do",
     "To produce high-quality, environmentally compliant bricks using VSK technology "
     "and biomass briquettes — providing affordable building materials that support "
     "Malawi's housing sector while protecting forests, creating jobs, and empowering "
     "communities in Mulanje District and beyond.",
     []),
    (DARK, "VISION", "Where We Are Going",
     "To build a scalable ecosystem for sustainable, affordable, and climate-resilient "
     "housing across Malawi and Southern Africa — positioning NEVAGS as the benchmark "
     "for eco-industrial construction innovation on the African continent.",
     []),
    (CARD, "OBJECTIVES", "Key Goals",
     "",
     ["Produce 4.8M+ bricks/year — zero firewood",
      "Scale to 200+ internal employees in 3 years",
      "Create 8,000+ jobs in the district ecosystem",
      "Achieve 30–45% annual ROI for investors",
      "Commission biogas 2nd kiln shaft in 14 months",
      "Expand to Blantyre, Lilongwe & Mzuzu markets",
      "Maintain ≥30% female employment — SDG 5",
      "Full ESG compliance — SDGs 3, 5, 13, 15"]),
]):
    x = Inches(0.5 + ci * 4.28)
    R(sl, x, COL_Y, Inches(4.1), COL_H, bg)
    R(sl, x, COL_Y, Inches(0.038), COL_H, ORANGE)
    T(sl, ttl,
      x+Inches(0.18), COL_Y+Inches(0.15), Inches(3.75), Inches(0.3),
      sz=9, bold=True, color=ORANGE)
    name_c = WHITE if ci < 2 else DARK
    T(sl, sub,
      x+Inches(0.18), COL_Y+Inches(0.45), Inches(3.75), Inches(0.42),
      sz=17, bold=True, color=name_c)
    R(sl, x+Inches(0.18), COL_Y+Inches(0.88), Inches(3.6), Inches(0.018),
      RGBColor(0x40,0x65,0x50) if ci < 2 else RGBColor(0xCC,0xCC,0xCC))
    if body:
        T(sl, body,
          x+Inches(0.18), COL_Y+Inches(1.0), Inches(3.75), Inches(4.1),
          sz=11, color=RGBColor(0xCC,0xEE,0xDD) if ci < 2 else GRAY)
    else:
        bx = sl.shapes.add_textbox(
            x+Inches(0.18), COL_Y+Inches(1.0), Inches(3.75), Inches(4.2))
        bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
        first = True
        for b in extras:
            if first: p = tf.paragraphs[0]; first = False
            else: p = tf.add_paragraph()
            p.space_before = Pt(6)
            r = p.add_run(); r.text = "  " + b
            r.font.size = Pt(10.5); r.font.color.rgb = GRAY; r.font.name = F

# Quote strip
R(sl, 0, H-Inches(0.48), W, Inches(0.48), DARK)
T(sl, '"Civilisations are remembered by what they build. '
  'The responsibility of our generation is to build differently."  — Chancy Tausi Tsonga',
  Inches(0.5), H-Inches(0.43), Inches(12.33), Inches(0.35),
  sz=9.5, italic=True, color=AMBER, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — PRODUCTS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Products & Services", "What We Produce", 6)

for ci, (bg, accent, tag, name, details, notes) in enumerate([
    (GREEN, ORANGE, "CORE PRODUCT",
     "Ordinary VSK Bricks",
     [("Production Cost", "K195 / brick"),
      ("QS Selling Price", "K300 / brick"),
      ("Gross Margin", "35%  ·  K105 / brick"),
      ("Annual Capacity", "3,000,000 bricks"),
      ("Annual Revenue",  "K900,000,000")],
     ["100% firewood-free — policy compliant",
      "Passes all Malawi govt standards",
      "Suitable for all construction types",
      "Competitive vs cement blocks"]),
    (DARK, ORANGE, "PREMIUM PRODUCT",
     "Face Bricks (Double-Faced)",
     [("Production Cost", "K335 / brick"),
      ("QS Selling Price", "K510–K650 / brick"),
      ("Gross Margin",    "34%  ·  K175 / brick"),
      ("Annual Capacity", "1,800,000 bricks"),
      ("Annual Revenue",  "K1,170,000,000")],
     ["Premium surface finish, no plastering",
      "Higher absolute margin per brick",
      "Commercial & high-end residential",
      "Supports early-phase revenue"]),
    (CARD, GREEN, "ECO INNOVATION",
     "Biomass Briquettes",
     [("Feedstock",      "Rice husks & agri waste"),
      ("Use",            "VSK kiln fuel (self-supply)"),
      ("Benefit",        "Replaces all firewood"),
      ("Community role", "Farmers paid per kg supply"),
      ("Future",         "Biogas 2nd shaft in 14 months")],
     ["Zero firewood — circular economy",
      "Turns farm waste into clean energy",
      "Community biomass supply chain",
      "Positions NEVAGS as biogas pioneer"]),
]):
    x = Inches(0.5 + ci * 4.28)
    H2 = Inches(5.42)
    R(sl, x, CONTENT_Y, Inches(4.1), H2, bg)
    R(sl, x, CONTENT_Y, Inches(0.038), H2, accent)
    T(sl, tag,
      x+Inches(0.18), CONTENT_Y+Inches(0.12), Inches(3.75), Inches(0.26),
      sz=8, bold=True, color=ORANGE)
    nc = WHITE if ci < 2 else DARK
    T(sl, name,
      x+Inches(0.18), CONTENT_Y+Inches(0.38), Inches(3.75), Inches(0.45),
      sz=16, bold=True, color=nc)
    R(sl, x+Inches(0.18), CONTENT_Y+Inches(0.84), Inches(3.6), Inches(0.018),
      RGBColor(0x40,0x65,0x50) if ci<2 else RGBColor(0xCC,0xCC,0xCC))
    dy = CONTENT_Y + Inches(0.92)
    for lbl, val in details:
        bx = sl.shapes.add_textbox(x+Inches(0.18), dy, Inches(3.7), Inches(0.32))
        bx.word_wrap = True; p = bx.text_frame.paragraphs[0]
        r1 = p.add_run(); r1.text = lbl + ":  "; r1.font.size = Pt(9)
        r1.font.bold = True
        r1.font.color.rgb = AMBER if ci < 2 else ORANGE; r1.font.name = F
        r2 = p.add_run(); r2.text = val; r2.font.size = Pt(9)
        r2.font.color.rgb = WHITE if ci < 2 else GRAY; r2.font.name = F
        dy += Inches(0.28)
    R(sl, x+Inches(0.18), dy, Inches(3.6), Inches(0.015),
      RGBColor(0x40,0x65,0x50) if ci < 2 else RGBColor(0xCC,0xCC,0xCC))
    dy += Inches(0.12)
    for b in notes:
        bx2 = sl.shapes.add_textbox(x+Inches(0.18), dy, Inches(3.7), Inches(0.28))
        p2 = bx2.text_frame.paragraphs[0]
        r = p2.add_run(); r.text = b; r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0xCC,0xEE,0xCC) if ci<2 else GRAY
        r.font.name = F; dy += Inches(0.28)

footer_bar(sl,
    "Annual Revenue: K900M (Ordinary)  +  K1,170M (Face Bricks)  =  K2.07 BILLION  "
    "·  Gross Profit: ~K774M (~USD 440,000)")


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — INVESTMENT OVERVIEW
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Investment Overview", "Capital Base, Returns & Financial Snapshot", 7)

# Capital stat row
for i, (num, lbl, sub, bg) in enumerate([
    ("K175M+",  "Founder's Equity",   "Charles Billy Nasala",   GREEN),
    ("EUR 175K","Atmosfair Loan",     "2nd Amendment, Apr 2026", DARK),
    ("K450M+",  "Total Capital",      "Founder + intl finance",  GREEN2),
    ("K2.07B",  "Revenue Potential",  "Annual MWK at capacity",  ORANGE),
]):
    stat_card(sl, Inches(0.5 + i*3.2), CONTENT_Y, Inches(3.05), Inches(1.4),
              num, lbl, sub, num_color=WHITE if i < 3 else WHITE, bg=bg)

# ROI row
for i, (num, lbl, body) in enumerate([
    ("30–45%", "Annual ROI Potential",
     "Based on full capacity at QS pricing. Policy-driven demand = long-term revenue."),
    ("3–5 yrs", "Investment Payback",
     "Conservative payback on total K450M+ investment. Expansion finance accelerates this."),
    ("$2.5M+", "Revenue — Year 3",
     "With scaled operations, revenue exceeds $2.5M USD within 3 years."),
    ("K774M",  "Annual Gross Profit",
     "~$440,000 USD annual gross profit at full production capacity."),
]):
    x = Inches(0.5 + i*3.2)
    R(sl, x, CONTENT_Y+Inches(1.52), Inches(3.05), Inches(1.58), CARD)
    T(sl, num, x+Inches(0.18), CONTENT_Y+Inches(1.6),
      Inches(2.7), Inches(0.55), sz=24, bold=True, color=GREEN)
    T(sl, lbl, x+Inches(0.18), CONTENT_Y+Inches(2.12),
      Inches(2.7), Inches(0.26), sz=9.5, bold=True, color=DARK)
    T(sl, body, x+Inches(0.18), CONTENT_Y+Inches(2.38),
      Inches(2.7), Inches(0.65), sz=8.5, color=GRAY)

# Revenue table
T(sl, "Annual Revenue & Profit Breakdown",
  Inches(0.5), CONTENT_Y+Inches(3.22), Inches(8), Inches(0.3),
  sz=12, bold=True, color=DARK)
table(sl, Inches(0.5), CONTENT_Y+Inches(3.55), Inches(12.33), 4, 5,
    [["Product",              "Capacity",          "Prod. Cost",    "QS Price",       "Annual Revenue"],
     ["Ordinary VSK Bricks",  "3,000,000 / year",  "K195 / brick",  "K300 / brick",   "K900,000,000"],
     ["Face Bricks (Premium)","1,800,000 / year",  "K335 / brick",  "K510–K650/brick","K1,170,000,000"],
     ["TOTAL",                "4,800,000 / year",  "—",             "—",              "K2,070,000,000"]],
    cw=[Inches(2.8), Inches(2.2), Inches(1.95), Inches(2.2), Inches(3.18)])


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — MARKET COMPARISON
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Pricing & Market Position",
             "Cost per m² of Wall — NEVAGS vs Alternatives", 8)

T(sl, "A standard half-brick wall (1 m²) requires 59 standard bricks or 12 cement blocks.",
  Inches(0.5), Inches(1.82), Inches(12.33), Inches(0.26),
  sz=10, italic=True, color=LGRAY)

BAR_MAX = Inches(7.2)
for i, (label, cost_str, pct, bg, note) in enumerate([
    ("Cement Blocks\n(12 blocks × K3,000)",
     "K36,000 / m²", 1.0, RGBColor(0xDC,0x26,0x26),
     "Most expensive option. Requires extra mortar, rebar and skilled labour. "
     "High transport cost to Mulanje District."),
    ("NEVAGS VSK Bricks\n(59 bricks × K300)",
     "K17,700 / m²", 0.492, GREEN,
     "LEGAL  ·  LOCAL  ·  AFFORDABLE — 51% cheaper than cement. "
     "Policy compliant. Eligible for all government and NGO projects."),
    ("Traditional Firewood Bricks\n(BANNED BY GOVERNMENT)",
     "K8,850 / m²", 0.246, RGBColor(0xCA,0x8A,0x04),
     "BANNED by Malawi Government. Using these bricks is illegal. "
     "Not eligible for any formal or government-funded construction project."),
]):
    y = CONTENT_Y + Inches(0.1) + Inches(i * 1.65)
    T(sl, label, Inches(0.5), y, Inches(3.1), Inches(0.6),
      sz=10, bold=True, color=DARK)
    bw = BAR_MAX * pct
    R(sl, Inches(3.7), y+Inches(0.06), bw, Inches(0.44), bg)
    T(sl, cost_str, Inches(3.7)+bw+Inches(0.12), y+Inches(0.1),
      Inches(1.8), Inches(0.35), sz=12, bold=True, color=DARK)
    T(sl, note, Inches(0.5), y+Inches(0.68), Inches(12.33), Inches(0.55),
      sz=8.5, color=GRAY)

R(sl, Inches(0.5), H-Inches(0.52), Inches(12.33), Inches(0.48), GREEN)
T(sl, "NEVAGS IS 51% CHEAPER THAN CEMENT BLOCKS  ·  "
  "The only legal, affordable, locally produced alternative in Mulanje District",
  Inches(0.7), H-Inches(0.48), Inches(11.9), Inches(0.40),
  sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — ENVIRONMENT & SDGs
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Environmental & Social Impact", "Built for a Sustainable Malawi", 9)

SDG_DATA = [
    (RGBColor(0x4C,0x9F,0x38), "SDG 3",
     "Good Health & Well-Being",
     ["Worker safety compliance ≥95%",
      "PPE provided for all 51 workers",
      "Reduced air & dust pollution",
      "Community health protection"]),
    (RGBColor(0xBF,0x0D,0x0D), "SDG 5",
     "Gender Equality",
     ["15 female employees (29%)",
      "Equal pay policy — no exceptions",
      "Female leadership roles",
      "Target: ≥30% women by Phase 3"]),
    (RGBColor(0x3F,0x7E,0x44), "SDG 13",
     "Climate Action",
     ["Zero firewood in all production",
      "Biomass briquette fuel only",
      "Carbon emission reduction",
      "Supports national climate targets"]),
    (RGBColor(0x56,0xC0,0x2B), "SDG 15",
     "Life on Land",
     ["Zero deforestation impact",
      "Protects the Mulanje Massif",
      "Biodiversity preservation",
      "Water catchment protection"]),
]
SW = Inches(2.9)
for i, (color, sdg, title, points) in enumerate(SDG_DATA):
    x = Inches(0.5 + i * 3.1)
    R(sl, x, CONTENT_Y, SW, Inches(0.7), color)
    T(sl, sdg, x+Inches(0.15), CONTENT_Y+Inches(0.08),
      SW-Inches(0.3), Inches(0.3), sz=20, bold=True, color=WHITE)
    T(sl, title, x+Inches(0.15), CONTENT_Y+Inches(0.38),
      SW-Inches(0.3), Inches(0.3), sz=9.5, color=WHITE)
    R(sl, x, CONTENT_Y+Inches(0.7), SW, Inches(3.72), CARD)
    bx = sl.shapes.add_textbox(
        x+Inches(0.15), CONTENT_Y+Inches(0.82), SW-Inches(0.2), Inches(3.5))
    bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
    first = True
    for pt in points:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run(); r.text = pt
        r.font.size = Pt(10.5); r.font.color.rgb = GRAY; r.font.name = F

# Zero waste strip
R(sl, Inches(0.5), CONTENT_Y+Inches(4.58), Inches(12.33), Inches(0.78), GREEN)
T(sl, "ZERO-WASTE CIRCULAR MODEL",
  Inches(0.7), CONTENT_Y+Inches(4.62), Inches(12), Inches(0.26),
  sz=9, bold=True, color=ORANGE)
T(sl, "Dust & clay waste → recovered and reused  "
  "·  Agricultural biomass → briquette kiln fuel  "
  "·  Organic waste → Biogas (2nd kiln, 14 months)  "
  "·  Kiln ash → soil amendment",
  Inches(0.7), CONTENT_Y+Inches(4.88), Inches(12), Inches(0.42),
  sz=11, color=RGBColor(0xCC,0xEE,0xCC), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — EMPLOYMENT & JOB CREATION
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Employment & Job Creation", "51 Jobs Today. 8,000+ in 3 Years.", 10)

# Four stat cards
for i, (num, lbl, sub, bg, nc) in enumerate([
    ("51",     "Current Employees", "15 Female · 36 Male",                GREEN, WHITE),
    ("200+",   "Internal Jobs",     "3-year target, all departments",     DARK,  WHITE),
    ("5,000+", "Community Moulders","Trained green-brick producers",      GREEN2,WHITE),
    ("2,500+", "Biomass Suppliers", "Rice husks & agri waste providers",  ORANGE,WHITE),
]):
    stat_card(sl, Inches(0.5+i*3.2), CONTENT_Y, Inches(3.05), Inches(1.4),
              num, lbl, sub, num_color=nc, bg=bg)

# Three ecosystem cards
for i, (ttl, body, accent) in enumerate([
    ("Community Brick Moulding Programme",
     "NEVAGS trains community members to produce VSK-quality green bricks at local level. "
     "We purchase directly from them and fire in our kiln. "
     "This creates genuine self-employment at community scale — no middleman.",
     GREEN),
    ("Biomass & Rice Husk Supply Chain",
     "We purchase rice husks, maize cobs and agricultural biomass from community farmers. "
     "Farm waste becomes clean energy income — turning zero-value materials into a "
     "predictable new rural revenue stream for Mulanje's agricultural sector.",
     ORANGE),
    ("Internal Growth — All Departments",
     "Beyond operations, NEVAGS scales across Sales, HR, Finance, Admin and a dedicated "
     "Mechanical Engineering & Maintenance department. "
     "200+ formal, permanent jobs within 3 years — above minimum wage, no exceptions.",
     GREEN2),
]):
    x = Inches(0.5 + i * 4.28)
    R(sl, x, CONTENT_Y+Inches(1.55), Inches(4.1), Inches(3.55), CARD)
    R(sl, x, CONTENT_Y+Inches(1.55), Inches(0.038), Inches(3.55), accent)
    T(sl, ttl,
      x+Inches(0.18), CONTENT_Y+Inches(1.68), Inches(3.8), Inches(0.38),
      sz=11, bold=True, color=DARK)
    T(sl, body,
      x+Inches(0.18), CONTENT_Y+Inches(2.1), Inches(3.8), Inches(2.8),
      sz=10, color=GRAY)

# Gender bar
R(sl, Inches(0.5), H-Inches(0.45), Inches(12.33), Inches(0.38), DARK)
R(sl, Inches(0.62), H-Inches(0.36), Inches(8.48), Inches(0.2), GREEN)
R(sl, Inches(9.1),  H-Inches(0.36), Inches(3.55), Inches(0.2), ORANGE)
T(sl, "Male  71%  (36 staff)",
  Inches(0.7), H-Inches(0.43), Inches(3.5), Inches(0.22),
  sz=8, bold=True, color=WHITE)
T(sl, "Female  29%  (15 staff)  —  SDG 5  ·  Target ≥30% by Phase 3",
  Inches(8.8), H-Inches(0.43), Inches(4.2), Inches(0.22),
  sz=8, bold=True, color=AMBER)


# ═══════════════════════════════════════════════════════════
# SLIDE 11 — WORKFORCE PLAN
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Workforce Plan", "KPI-Driven Phased Scaling — Apr to Jul 2026", 11)

table(sl, Inches(0.5), CONTENT_Y, Inches(12.33), 6, 5,
    [["Phase",              "Period",           "Staff",            "Daily Output Target",    "Key KPIs"],
     ["Startup",            "Wk 1–2  Apr 6–19", "18 staff",         "Site commissioning",     "Safety ≥95%  ·  Downtime <10%"],
     ["Phase 1",            "Wk 3–4  Apr–May",  "33 staff",         "≥3,000 bricks/day",  "Output ≥3,000/day  ·  Cost ≤K195/brick"],
     ["Phase 2",            "Wk 5–8  May",      "39–40 staff",      "≥5,000 bricks/day",  "Kiln util >80%  ·  Orders >50,000/wk"],
     ["Phase 3 (Full)",     "Wk 9–13  Jun–Jul", "45 staff",         "≥7,000 bricks/day",  "Active clients ≥5  ·  Double-shift"],
     ["Expansion",          "2026–2029",         "200+ internal",    "Doubled capacity",       "8,000+ ecosystem jobs  ·  Biogas kiln"]],
    cw=[Inches(1.7), Inches(2.2), Inches(1.65), Inches(3.0), Inches(3.78)])

T(sl, "Compensation Philosophy",
  Inches(0.5), CONTENT_Y+Inches(2.65), Inches(12), Inches(0.3),
  sz=12, bold=True, color=DARK)

for i, (ttl, body, bg) in enumerate([
    ("Senior & Management Staff",
     "Competitive market-rate remuneration — "
     "above-market base salaries reflecting expertise and business contribution.",
     GREEN),
    ("General Workers & Supervisors",
     "Above statutory minimum wage — no exceptions. Stable year-round income "
     "with career progression tied directly to measurable KPI performance.",
     DARK),
    ("Community Moulders & Biomass Suppliers",
     "Direct purchase payments at guaranteed rates. Offtake agreements create "
     "predictable, sustainable income for community entrepreneurs.",
     GREEN2),
]):
    x = Inches(0.5 + i * 4.28)
    R(sl, x, CONTENT_Y+Inches(3.05), Inches(4.1), Inches(2.1), bg)
    R(sl, x, CONTENT_Y+Inches(3.05), Inches(0.038), Inches(2.1), ORANGE)
    T(sl, ttl,
      x+Inches(0.18), CONTENT_Y+Inches(3.15), Inches(3.8), Inches(0.35),
      sz=10, bold=True, color=ORANGE)
    T(sl, body,
      x+Inches(0.18), CONTENT_Y+Inches(3.52), Inches(3.8), Inches(1.5),
      sz=10, color=RGBColor(0xCC,0xDD,0xCC))


# ═══════════════════════════════════════════════════════════
# SLIDE 12 — 3-MONTH BUDGET
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Operational Budget", "3-Month Operational Budget Breakdown", 12)

budget_items = [
    ("Salaries & Wages",       "K20.3M",  "All 51 staff, 3 months — above minimum wage"),
    ("Raw Materials",           "K7.5M",   "Clay, binding agents, consumables"),
    ("Machinery & Tools",       "K4.6M",   "Maintenance, spare parts, equipment"),
    ("PPE & Safety",            "K4.7M",   "Personal protective equipment for all staff"),
    ("Office & IT",             "K2.6M",   "Admin, communications, office operations"),
    ("Utilities & Energy",      "K1.6M",   "ESCOM power, water, fuel logistics"),
    ("Marketing & Sales",       "K1.0M",   "Client outreach, DEC, tenders"),
    ("Insurance & Admin",       "K0.7M",   "Business insurance, legal, compliance"),
]

# Left: table
table(sl, Inches(0.5), CONTENT_Y, Inches(7.6), 9, 3,
    [["Budget Line",       "Amount (MWK)", "Description"]] +
    [[i, a, d] for i, a, d in budget_items] ,
    cw=[Inches(2.5), Inches(1.5), Inches(3.6)])

# Right: summary cards + total
R(sl, Inches(8.3), CONTENT_Y, Inches(4.53), Inches(1.58), GREEN)
R(sl, Inches(8.3), CONTENT_Y, Inches(0.038), Inches(1.58), ORANGE)
T(sl, "TOTAL 3-MONTH BUDGET",
  Inches(8.5), CONTENT_Y+Inches(0.12), Inches(4.1), Inches(0.3),
  sz=9, bold=True, color=ORANGE)
T(sl, "K43.0 Million",
  Inches(8.5), CONTENT_Y+Inches(0.42), Inches(4.1), Inches(0.7),
  sz=32, bold=True, color=WHITE)
T(sl, "Approx. USD 24,700  ·  3-month operational runway",
  Inches(8.5), CONTENT_Y+Inches(1.1), Inches(4.1), Inches(0.32),
  sz=9, color=RGBColor(0xCC,0xDD,0xCC))

# Bar chart (visual approximation)
T(sl, "Budget Distribution",
  Inches(8.3), CONTENT_Y+Inches(1.75), Inches(4.53), Inches(0.28),
  sz=10, bold=True, color=DARK)

bar_colors = [GREEN, ORANGE, GREEN2, RGBColor(0xCA,0x8A,0x04),
              DARK, GREEN3, ORANGE, RGBColor(0x6B,0x7B,0x8D)]
bar_values = [20.3, 7.5, 4.6, 4.7, 2.6, 1.6, 1.0, 0.7]
bar_max = 20.3
bar_total_w = Inches(4.3)
for bi, (lbl, val) in enumerate(zip([x[0] for x in budget_items], bar_values)):
    by = CONTENT_Y + Inches(2.1 + bi * 0.58)
    bw = bar_total_w * (val / bar_max)
    short_lbl = lbl.split(" ")[0]
    T(sl, short_lbl, Inches(8.3), by, Inches(1.55), Inches(0.34),
      sz=8.5, color=GRAY)
    R(sl, Inches(9.85), by+Inches(0.04), bw, Inches(0.26), bar_colors[bi])
    T(sl, f"K{val}M", Inches(9.85)+bw+Inches(0.08), by, Inches(0.8), Inches(0.3),
      sz=8, color=DARK)


# ═══════════════════════════════════════════════════════════
# SLIDE 13 — PARTNERS & STAKEHOLDERS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Partners & Stakeholders", "Backed by Leading Institutions", 13)

partners = [
    ("GIZ",      "Deutsche Gesellschaft fur Internationale Zusammenarbeit",
     "Technical & Financial Partner",   "German Development Cooperation"),
    ("Atmosfair", "German Climate Protection NGO  ·  Berlin",
     "Carbon Finance & Loan Provider",  "EUR 175,000 (2nd Amendment, Apr 2026)"),
    ("CCODE",    "Centre for Community Organisation & Development",
     "Community Engagement",            "Local Development Partner"),
    ("TERA",     "Technical & Environmental Research Associates",
     "Environmental Compliance",        "Research & Advisory Partner"),
    ("MUBAS",    "Malawi University of Business & Applied Sciences",
     "Academic Research",               "Skills Development & Knowledge"),
    ("CIRA",     "Construction Industry Regulatory Authority",
     "Industry Compliance",             "Market Standards & Development"),
]
for i, (abbr, org, role, sub) in enumerate(partners):
    col = i % 3; row = i // 3
    x = Inches(0.5 + col * 4.28)
    y = CONTENT_Y + Inches(row * 2.58)
    acc = ORANGE if row == 0 else GREEN
    R(sl, x, y, Inches(4.1), Inches(2.42), CARD)
    R(sl, x, y, Inches(0.038), Inches(2.42), acc)
    T(sl, abbr,
      x+Inches(0.18), y+Inches(0.12), Inches(3.8), Inches(0.55),
      sz=24, bold=True, color=GREEN)
    T(sl, org,
      x+Inches(0.18), y+Inches(0.68), Inches(3.8), Inches(0.38),
      sz=8.5, color=GRAY)
    T(sl, role,
      x+Inches(0.18), y+Inches(1.06), Inches(3.8), Inches(0.3),
      sz=10, bold=True, color=DARK)
    T(sl, sub,
      x+Inches(0.18), y+Inches(1.36), Inches(3.8), Inches(0.85),
      sz=9, color=LGRAY)

footer_bar(sl,
    "Government  ·  NGOs & Development Finance  ·  "
    "Academic Institutions  ·  Industry Regulators  ·  Community Organisations")


# ═══════════════════════════════════════════════════════════
# SLIDE 14 — IMPLEMENTATION ROADMAP
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "Implementation Roadmap", "Phased Growth Plan — 2026 to 2029", 14)

phases = [
    (GREEN,  "PHASE 1", "Foundation &\nCommissioning", "Apr – May 2026",
     True,    # show commissioning badge
     ["VSK kiln commissioning complete",
      "18 → 33 staff on-boarded",
      "Moulding operations launched",
      "≥3,000 bricks/day output",
      "First client engagements"]),
    (ORANGE, "PHASE 2", "Scale\nProduction", "May – Jun 2026",
     False,
     ["Output scaled to ≥5,000/day",
      "39–40 staff deployed",
      "Kiln utilisation >80%",
      "50,000+ bricks/week orders",
      "Market penetration active"]),
    (GREEN2, "PHASE 3", "Full-Scale\nOperations", "Jun – Jul 2026",
     False,
     ["45 staff at full capacity",
      "Output ≥7,000 bricks/day",
      "≥5 active bulk clients",
      "Double-shift introduced",
      "Revenue targets met"]),
    (RGBColor(0x3B,0x82,0xF6), "PHASE 4", "Expansion &\nInnovation", "2026 – 2029",
     False,
     ["200+ internal staff",
      "8,000+ ecosystem jobs",
      "Biogas 2nd kiln shaft",
      "Revenue >$2.5M USD",
      "Southern Africa expansion"]),
]
PW = Inches(2.9)
for i, (bg, ph, title, period, show_badge, pts) in enumerate(phases):
    x = Inches(0.5 + i * 3.21)
    # Phase label header
    R(sl, x, CONTENT_Y, PW, Inches(0.58), bg)
    T(sl, ph, x+Inches(0.15), CONTENT_Y+Inches(0.06),
      PW-Inches(0.3), Inches(0.28), sz=10, bold=True, color=WHITE)
    T(sl, period, x+Inches(0.15), CONTENT_Y+Inches(0.32),
      PW-Inches(0.3), Inches(0.22), sz=8.5, color=WHITE)
    # Body panel
    R(sl, x, CONTENT_Y+Inches(0.58), PW, Inches(4.58), CARD)
    T(sl, title, x+Inches(0.15), CONTENT_Y+Inches(0.65),
      PW-Inches(0.3), Inches(0.55), sz=14, bold=True, color=DARK)
    # Commissioning badge in Phase 1
    if show_badge:
        R(sl, x+Inches(0.15), CONTENT_Y+Inches(1.22), PW-Inches(0.25), Inches(0.36),
          RGBColor(0xFF,0xF3,0xE0))
        T(sl, "27 May 2026 — Commissioning Ceremony",
          x+Inches(0.22), CONTENT_Y+Inches(1.25), PW-Inches(0.38), Inches(0.28),
          sz=8.5, bold=True, color=ORANGE)
        bullet_start_y = CONTENT_Y + Inches(1.65)
    else:
        bullet_start_y = CONTENT_Y + Inches(1.28)
    bx = sl.shapes.add_textbox(
        x+Inches(0.15), bullet_start_y, PW-Inches(0.2), Inches(3.4))
    bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
    first = True
    for pt in pts:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        p.space_before = Pt(7)
        r = p.add_run(); r.text = pt
        r.font.size = Pt(10); r.font.color.rgb = GRAY; r.font.name = F

# Biogas spotlight
R(sl, Inches(0.5), H-Inches(0.55), Inches(12.33), Inches(0.5), GREEN)
R(sl, Inches(0.5), H-Inches(0.55), Inches(0.038), Inches(0.5), ORANGE)
T(sl, "IN 14 MONTHS: Biogas-Powered 2nd VSK Kiln Shaft — "
  "First Biogas Brick Producer in Malawi  ·  Doubles Capacity",
  Inches(0.7), H-Inches(0.5), Inches(11.9), Inches(0.42),
  sz=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 15 — WHY DISTRICT / WHY NEVAGS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
section_head(sl, "A Relationship of Mutual Benefit",
             "Why Mulanje Needs NEVAGS  ·  Why NEVAGS Needs Mulanje", 15)

for ci, (bg, accent, heading, pairs) in enumerate([
    (GREEN, ORANGE, "WHY MULANJE DISTRICT NEEDS NEVAGS", [
        ("Policy Compliance",
         "NEVAGS is the only VSK producer in the district — the only legal solution "
         "to the firewood ban. Every government and NGO project needs compliant bricks."),
        ("8,000+ Jobs",
         "The largest potential job creator in Mulanje history — "
         "from factory floor to the community supply chain, within 3 years."),
        ("Forest Protection",
         "Every NEVAGS brick protects the Mulanje Massif's forests, "
         "water catchments and biodiversity."),
        ("Black-Owned Pioneer",
         "One of the first black-owned industrial enterprises in Mulanje. "
         "Supporting NEVAGS sends a historic signal for the district."),
        ("Local Economy",
         "Revenue stays in Mulanje — wages, local supply purchases, "
         "and growing district tax revenues."),
    ]),
    (DARK, ORANGE, "WHY NEVAGS NEEDS MULANJE DISTRICT", [
        ("DEC Endorsement",
         "Formal endorsement unlocks government procurement — "
         "millions of bricks in guaranteed long-term contracts."),
        ("Workforce",
         "Mulanje's community is our talent pool. District cooperation "
         "enables the community training and moulding programme."),
        ("Biomass Supply",
         "Rice husks and farm biomass from Mulanje's farmers are our kiln fuel. "
         "We need the agricultural sector as a strategic partner."),
        ("Infrastructure",
         "Stable ESCOM power, road access and land tenure security "
         "require district-level facilitation."),
        ("Policy Champion",
         "As a pioneering black enterprise, NEVAGS needs the DEC "
         "to champion, protect and amplify our growth."),
    ]),
]):
    x = Inches(0.5 + ci * 6.5)
    R(sl, x, CONTENT_Y, Inches(6.2), Inches(5.42), bg)
    R(sl, x, CONTENT_Y, Inches(0.038), Inches(5.42), accent)
    T(sl, heading,
      x+Inches(0.18), CONTENT_Y+Inches(0.12), Inches(5.9), Inches(0.28),
      sz=9, bold=True, color=ORANGE)
    R(sl, x+Inches(0.18), CONTENT_Y+Inches(0.44), Inches(5.9), Inches(0.018),
      RGBColor(0x40,0x65,0x50))
    bx = sl.shapes.add_textbox(
        x+Inches(0.18), CONTENT_Y+Inches(0.52), Inches(5.85), Inches(4.7))
    bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
    first = True
    for title_t, body_t in pairs:
        if not first:
            p = tf.add_paragraph(); p.space_before = Pt(2)
            r = p.add_run(); r.text = ""; r.font.name = F
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(6)
        r = p.add_run(); r.text = title_t
        r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = AMBER; r.font.name = F
        p2 = tf.add_paragraph(); p2.space_before = Pt(2)
        r2 = p2.add_run(); r2.text = body_t
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(0xCC,0xEE,0xCC) if ci == 0 else RGBColor(0xCC,0xCC,0xCC)
        r2.font.name = F


# ═══════════════════════════════════════════════════════════
# SLIDE 16 — THE ASK & CONTACT
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
R(sl, 0, 0, W, H, WHITE)
# Dark left panel
R(sl, 0, 0, Inches(0.5), H, GREEN)
# Orange accent top
R(sl, Inches(0.5), 0, Inches(12.83), Inches(0.06), ORANGE)

page_header(sl, "Partnership Invitation", 16, invert=False)

T(sl, "PARTNERSHIP INVITATION",
  Inches(0.65), Inches(0.7), Inches(12), Inches(0.28),
  sz=9, bold=True, color=ORANGE)
T(sl, "Let’s Build Malawi’s\nSustainable Future Together",
  Inches(0.65), Inches(0.95), Inches(10), Inches(1.1),
  sz=34, bold=True, color=DARK)
T(sl, "NEVAGS asks the District Executive Committee for three things:",
  Inches(0.65), Inches(2.1), Inches(12), Inches(0.28),
  sz=11, color=GRAY)

# 3 ask cards
for i, (num, title, body) in enumerate([
    ("1", "Recognition & Endorsement",
     "Formally recognise NEVAGS as a strategic district partner — opening "
     "government procurement, NGO supply contracts and public sector projects."),
    ("2", "Community Programme Support",
     "Facilitate community outreach, training endorsement and biomass supply chain "
     "development through district channels and networks."),
    ("3", "Infrastructure Facilitation",
     "Support stable ESCOM power, road infrastructure and land tenure security — "
     "the foundational requirements for NEVAGS to scale and deliver on all commitments."),
]):
    x = Inches(0.65 + i * 4.2)
    R(sl, x, Inches(2.52), Inches(4.0), Inches(2.32), CARD)
    R(sl, x, Inches(2.52), Inches(0.038), Inches(2.32), GREEN if i != 1 else ORANGE)
    T(sl, num, x+Inches(0.15), Inches(2.6), Inches(0.45), Inches(0.45),
      sz=24, bold=True, color=RGBColor(0xD1,0xD5,0xDB))
    T(sl, title, x+Inches(0.65), Inches(2.6), Inches(3.2), Inches(0.42),
      sz=12, bold=True, color=DARK)
    T(sl, body, x+Inches(0.18), Inches(3.06), Inches(3.7), Inches(1.65),
      sz=9.5, color=GRAY)

# Contact cards
T(sl, "CONTACTS",
  Inches(0.65), Inches(5.0), Inches(12), Inches(0.26),
  sz=9, bold=True, color=ORANGE)

for i, (role, name, lines) in enumerate([
    ("Managing Director",
     "Charles Billy Nasala",
     ["+265 888 34 75 75  ·  +265 99 751 0160",
      "nasalacharles.b@gmail.com"]),
    ("Founding Engineer & BD Manager",
     "Chancy Tausi Tsonga",
     ["+265 984 000 366  ·  WA: +27 764 998 4601",
      "chancy.tsonga@yahoo.com  ·  chancytsonga.com"]),
    ("Careers & Head Office",
     "NEVAGS Human Resources",
     ["careers.nevags@gmail.com",
      "Musewu, Mulanje District  ·  P.O. Box 90"]),
]):
    x = Inches(0.65 + i * 4.2)
    R(sl, x, Inches(5.28), Inches(4.0), Inches(1.88), CARD)
    R(sl, x, Inches(5.28), Inches(0.038), Inches(1.88), GREEN)
    T(sl, role, x+Inches(0.18), Inches(5.36), Inches(3.7), Inches(0.24),
      sz=8.5, bold=True, color=ORANGE)
    T(sl, name, x+Inches(0.18), Inches(5.60), Inches(3.7), Inches(0.36),
      sz=13, bold=True, color=DARK)
    for j, ln in enumerate(lines):
        T(sl, ln, x+Inches(0.18), Inches(5.98 + j*0.3), Inches(3.7), Inches(0.28),
          sz=9, color=GRAY)

# Final footer
R(sl, Inches(0.5), H-Inches(0.38), Inches(12.33), Inches(0.38), GREEN)
T(sl, "NEVAGS ECO BRICK & CONSTRUCTION  "
  "·  Reg. 46289  "
  "·  Musewu, Mulanje, Malawi  "
  "·  Building Tomorrow Sustainably",
  Inches(0.7), H-Inches(0.34), Inches(11.9), Inches(0.28),
  sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
out = "NEVAGS_DEC_Presentation_2026.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
